from pathlib import Path
from datetime import datetime


class OutputGeneratorAgent:

    def _create_pdf(self, plan: dict, output_dir: Path) -> Path:
        """Generate PDF with cover, stops table, day plan list, cost summary."""
        from fpdf import FPDF

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Cover page
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 24)
        pdf.cell(0, 20, "Reiseplan", ln=True, align="C")
        pdf.set_font("Helvetica", "", 14)
        pdf.cell(0, 10, f"Von: {plan.get('start_location', '')}", ln=True, align="C")

        stops = plan.get("stops", [])
        if stops:
            last = stops[-1]
            pdf.cell(0, 10, f"Nach: {last.get('region', '')} ({last.get('country', '')})", ln=True, align="C")

        pdf.cell(0, 10, f"Erstellt: {datetime.now().strftime('%d.%m.%Y')}", ln=True, align="C")

        # Stops table
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, "Reise-Stopps", ln=True)
        pdf.ln(5)

        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(40, 8, "Stop", border=1)
        pdf.cell(30, 8, "Land", border=1)
        pdf.cell(25, 8, "Tag", border=1)
        pdf.cell(20, 8, "Nächte", border=1)
        pdf.cell(75, 8, "Unterkunft", border=1)
        pdf.ln()

        pdf.set_font("Helvetica", "", 9)
        for stop in stops:
            acc = stop.get("accommodation", {}) or {}
            acc_name = acc.get("name", "-") if isinstance(acc, dict) else "-"
            pdf.cell(40, 7, stop.get("region", "")[:20], border=1)
            pdf.cell(30, 7, stop.get("country", ""), border=1)
            pdf.cell(25, 7, str(stop.get("arrival_day", "")), border=1)
            pdf.cell(20, 7, str(stop.get("nights", "")), border=1)
            pdf.cell(75, 7, acc_name[:35], border=1)
            pdf.ln()

        # Day plan
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, "Tagesplan", ln=True)
        pdf.ln(5)

        for dp in plan.get("day_plans", []):
            pdf.set_font("Helvetica", "B", 11)
            day_str = f"Tag {dp.get('day', '')}: {dp.get('title', '')}"
            pdf.multi_cell(0, 8, day_str)
            pdf.set_font("Helvetica", "", 9)
            pdf.multi_cell(0, 6, dp.get("description", ""))
            pdf.ln(3)

        # Cost summary
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, "Kostenübersicht", ln=True)
        pdf.ln(5)

        cost = plan.get("cost_estimate", {}) or {}
        pdf.set_font("Helvetica", "", 11)
        items = [
            ("Unterkunft", cost.get("accommodations_chf", 0)),
            ("Fähren", cost.get("ferries_chf", 0)),
            ("Aktivitäten", cost.get("activities_chf", 0)),
            ("Verpflegung", cost.get("food_chf", 0)),
            ("Treibstoff", cost.get("fuel_chf", 0)),
        ]
        for label, amount in items:
            pdf.cell(100, 8, label)
            pdf.cell(0, 8, f"CHF {amount:,.0f}", ln=True)

        pdf.ln(5)
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(100, 8, "Total")
        pdf.cell(0, 8, f"CHF {cost.get('total_chf', 0):,.0f}", ln=True)
        pdf.cell(100, 8, "Verbleibendes Budget")
        pdf.cell(0, 8, f"CHF {cost.get('budget_remaining_chf', 0):,.0f}", ln=True)

        # Save
        filename = f"reiseplan_{plan.get('job_id', 'unknown')}.pdf"
        output_path = output_dir / filename
        pdf.output(str(output_path))
        return output_path

    def _create_pptx(self, plan: dict, output_dir: Path) -> Path:
        """Generate PPTX with title slide + per-stop slides + cost summary."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # blank
        title_layout = prs.slide_layouts[0]  # title slide

        # Title slide
        slide = prs.slides.add_slide(title_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        stops = plan.get("stops", [])
        title.text = "Reiseplan"
        dest = stops[-1].get("region", "") if stops else ""
        subtitle.text = f"Von {plan.get('start_location', '')} nach {dest}\n{datetime.now().strftime('%d.%m.%Y')}"

        # Stop slides
        content_layout = prs.slide_layouts[1]
        for stop in stops:
            slide = prs.slides.add_slide(content_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = f"{stop.get('region', '')} ({stop.get('country', '')})"

            acc = stop.get("accommodation", {}) or {}
            acc_name = acc.get("name", "-") if isinstance(acc, dict) else "-"
            acc_price = acc.get("total_price_chf", 0) if isinstance(acc, dict) else 0

            acts = stop.get("top_activities", [])
            act_names = "\n".join(f"  • {a.get('name', '')}" for a in acts[:3])

            rests = stop.get("restaurants", [])
            rest_names = "\n".join(f"  • {r.get('name', '')}" for r in rests[:2])

            tf = body.text_frame
            tf.text = f"Ankunftstag: {stop.get('arrival_day', '')}\nNächte: {stop.get('nights', '')}\n"
            tf.text += f"\nUnterkunft: {acc_name} (CHF {acc_price:,.0f})\n"
            if act_names:
                tf.text += f"\nAktivitäten:\n{act_names}\n"
            if rest_names:
                tf.text += f"\nRestaurants:\n{rest_names}\n"

        # Cost summary slide
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Kostenübersicht"

        cost = plan.get("cost_estimate", {}) or {}
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = (
            f"Unterkunft:     CHF {cost.get('accommodations_chf', 0):,.0f}\n"
            f"Aktivitäten:    CHF {cost.get('activities_chf', 0):,.0f}\n"
            f"Verpflegung:    CHF {cost.get('food_chf', 0):,.0f}\n"
            f"Treibstoff:     CHF {cost.get('fuel_chf', 0):,.0f}\n"
            f"Fähren:         CHF {cost.get('ferries_chf', 0):,.0f}\n"
            f"─────────────────────────\n"
            f"Total:          CHF {cost.get('total_chf', 0):,.0f}\n"
            f"Rest-Budget:    CHF {cost.get('budget_remaining_chf', 0):,.0f}"
        )

        filename = f"reiseplan_{plan.get('job_id', 'unknown')}.pptx"
        output_path = output_dir / filename
        prs.save(str(output_path))
        return output_path
